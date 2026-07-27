#!/usr/bin/env python3
"""
build-pptx.py — Assemble the 14 rendered slide MP4s into a single PowerPoint
where each video is full-screen and set to PLAY AUTOMATICALLY on slide entry,
with a fade transition between slides.

Present it, talk over each slide, and click to advance — the clip on the next
slide starts playing by itself.

Usage:
    python3 scripts/build-pptx.py            # reads ppt/slide-01..14.mp4
Output:
    ppt/Capital-Command-Manifesto.pptx
"""
import os
import subprocess
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPT_DIR = os.path.join(HERE, "ppt")
FFMPEG = os.path.join(
    HERE, "node_modules", "@remotion", "compositor-linux-x64-gnu", "ffmpeg"
)
if not os.path.exists(FFMPEG):
    FFMPEG = "ffmpeg"

# Slide length in seconds — mirrors DURATION_SECONDS in src/config.ts.
DURATIONS = {
    1: 9, 2: 8, 3: 7, 4: 7, 5: 9, 6: 7, 7: 7,
    8: 8, 9: 7, 10: 8, 11: 9, 12: 8, 13: 8, 14: 11,
}

TITLES = {
    1: "The Hook — 1995 → 2026",
    2: "The Resistance",
    3: "The Fear",
    4: "It Rewrites It",
    5: "The Precedent",
    6: "The Claim",
    7: "The Consequence",
    8: "The Choice",
    9: "The Pivot",
    10: "One Person",
    11: "CoLateral — Product One",
    12: "The Vision",
    13: "The Philosophy",
    14: "The Empire",
}

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def poster_for(mp4_path, out_png):
    """Grab the last frame of the clip as the slide's static poster image."""
    subprocess.run(
        [FFMPEG, "-y", "-sseof", "-0.1", "-i", mp4_path, "-vframes", "1", out_png],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def autoplay_timing_xml(spid, duration_ms):
    """
    The exact timing tree PowerPoint writes for a video set to
    'Start: Automatically'. The video plays as part of the main sequence the
    moment the slide appears.
    """
    return f"""<p:timing xmlns:p="{P_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="{duration_ms}" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:video>
            <p:cMediaNode vol="80000">
              <p:cTn id="7" fill="hold" display="0">
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


def fade_transition_xml():
    """A smooth fade transition, advancing on click (so you control the pace)."""
    return (
        f'<p:transition xmlns:p="{P_NS}" spd="med">'
        f"<p:fade thruBlk=\"1\"/>"
        f"</p:transition>"
    )


def set_autoplay_and_transition(slide, movie, duration_s):
    el = slide._element
    spid = movie._element.find(qn("p:nvPicPr") + "/" + qn("p:cNvPr")).get("id")

    # Replace the default (play-on-click) timing with the autoplay tree.
    old = el.find(qn("p:timing"))
    if old is not None:
        el.remove(old)
    new_timing = etree.fromstring(autoplay_timing_xml(spid, int(duration_s * 1000)))

    # Schema order in CT_Slide: cSld, clrMapOvr, transition, timing, extLst.
    # Insert transition before timing.
    transition = etree.fromstring(fade_transition_xml())
    clr = el.find(qn("p:clrMapOvr"))
    if clr is not None:
        clr.addnext(transition)
    else:
        el.insert(0, transition)
    transition.addnext(new_timing)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i in range(1, 15):
        mp4 = os.path.join(PPT_DIR, f"slide-{i:02d}.mp4")
        if not os.path.exists(mp4):
            raise FileNotFoundError(f"Missing render: {mp4}")
        poster = os.path.join(PPT_DIR, f"poster-{i:02d}.png")
        poster_for(mp4, poster)

        slide = prs.slides.add_slide(blank)
        movie = slide.shapes.add_movie(
            mp4,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
            poster_frame_image=poster,
            mime_type="video/mp4",
        )
        set_autoplay_and_transition(slide, movie, DURATIONS[i])
        print(f"  + slide {i:02d}  {TITLES[i]}")

    out = os.path.join(PPT_DIR, "Capital-Command-Manifesto.pptx")
    prs.save(out)
    size_mb = os.path.getsize(out) / (1024 * 1024)
    print(f"\n✅ Saved {out}  ({size_mb:.1f} MB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
