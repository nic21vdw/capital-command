#!/usr/bin/env python3
"""
Build a PowerPoint (.pptx) that embeds the rendered Remotion intro and plays it
AUTOMATICALLY when you land on its slide (no click needed).

python-pptx's add_movie() always wires a video to "play on click". To make it
autoplay on slide entry we replace the slide's <p:timing> tree with one whose
media trigger uses nodeType="afterEffect" + <p:cond delay="0"/> — this is the
same shape PowerPoint writes when you pick Playback ▸ Start ▸ Automatically.

Run:  python3 build_pptx.py
Out:  out/free-ai-builds-intro.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
VIDEO = os.path.join(HERE, "out", "free-ai-builds-intro.mp4")
POSTER = os.path.join(HERE, "out", "poster.png")
OUT = os.path.join(HERE, "out", "free-ai-builds-intro.pptx")

# Portrait 9:16 to match the vertical 1080x1920 video exactly (full-bleed).
SLIDE_W = Emu(6858000)   # 7.5 in
SLIDE_H = Emu(12192000)  # 13.333 in

# "Signal" theme
BG = RGBColor(0x0F, 0x0F, 0x0F)
ACCENT = RGBColor(0xFF, 0x4B, 0x1F)
HIGHLIGHT = RGBColor(0xFF, 0xD2, 0x3F)
TEXT = RGBColor(0xFA, 0xFA, 0xFA)

VIDEO_MS = 23000  # 690 frames / 30 fps


def solid_bg(slide, color):
    """Paint an entire slide a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, spacing) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Poppins"
        if spacing:
            _letter_spacing(r, spacing)
    return tb


def _letter_spacing(run, pts):
    """Set character spacing (spc, in 1/100 pt) on a run."""
    run.font._rPr.set("spc", str(int(pts * 100)))


def autoplay_timing_xml(shape_id):
    """
    <p:timing> that starts the media automatically on slide entry.
    Key vs. click-to-play: nodeType='afterEffect' and <p:cond delay='0'/>.
    """
    return f"""
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetClass="mediacall" presetID="1" nodeType="afterEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="{VIDEO_MS}" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:video>
            <p:cMediaNode vol="80000">
              <p:cTn id="7" fill="hold" display="0">
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                <p:endCondLst>
                  <p:cond evt="onStopAudio" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
                </p:endCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>""".strip()


def set_autoplay(slide, shape_id):
    """Attach the autoplay <p:timing> to the slide, replacing any existing one."""
    sld = slide._element
    for existing in sld.findall(qn("p:timing")):
        sld.remove(existing)
    timing = etree.fromstring(autoplay_timing_xml(shape_id))
    sld.append(timing)


def main():
    for f in (VIDEO, POSTER):
        if not os.path.exists(f):
            raise SystemExit(f"Missing input: {f} — render the video first.")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # fully blank layout

    # ---- Slide 1: title / lead-in (so advancing INTO slide 2 triggers autoplay)
    s1 = prs.slides.add_slide(blank)
    solid_bg(s1, BG)
    add_text(
        s1, Emu(600000), Emu(4600000), SLIDE_W - Emu(1200000), Emu(3000000),
        [
            ("FREE AI BUILDS", 54, TEXT, True, 0),
            ("A NEW SERIES", 24, ACCENT, True, 6),
            ("", 12, TEXT, False, 0),
            ("Advance to the next slide ▸", 18, HIGHLIGHT, False, 0),
        ],
    )

    # ---- Slide 2: full-bleed video that autoplays on entry
    s2 = prs.slides.add_slide(blank)
    solid_bg(s2, BG)
    movie = s2.shapes.add_movie(
        VIDEO, 0, 0, SLIDE_W, SLIDE_H,
        poster_frame_image=POSTER, mime_type="video/mp4",
    )
    set_autoplay(s2, movie.shape_id)

    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"  slide 1: title lead-in")
    print(f"  slide 2: video autoplays on entry (shape id {movie.shape_id})")


if __name__ == "__main__":
    main()
