"""Build a click-to-advance PowerPoint from the rendered Fable 5 segments.

One video per slide, full-bleed, embedded (not linked) so the .pptx is
self-contained. In PowerPoint: click the play button to start a segment,
click to advance to the next slide when you're done talking.
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

# True  -> each video auto-plays the instant you land on its slide; you only
#          click to advance to the next slide.
# False -> click the video to start it (manual play).
AUTOPLAY = True

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "out")

# (video file, poster image, speaker-note cue)
SEGMENTS = [
    ("oc-01-intro.mp4", "p-01.png",
     "Fable 5 just came back - and if you write code, the way you pay for it quietly changed today."),
    ("oc-02-timeline.mp4", "p-02.png",
     "June 9 it goes public, June 12 it's pulled under export controls, July 1 the controls lift and it goes worldwide."),
    ("oc-03-coding.mp4", "p-03.png",
     "On SWE-Bench Pro it's around 80%. Stripe said it compressed months of engineering into days."),
    ("oc-04-gpt-pressure.mp4", "p-04.png",
     "GPT-5.6 is still limited preview, vetted orgs only. Fable 5 is just available to everyone. Access beats a benchmark you can't log into."),
    ("oc-05-coder-catch.mp4", "p-05.png",
     "The weekly limits were actually reset - twice. The real change: Fable 5 now leaves your subscription limits and moves to usage credits."),
    ("oc-06-price.mp4", "p-06.png",
     "That's $50 per million output tokens. If you run big agentic coding jobs, plan for it."),
]

prs = Presentation()
# 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # fully blank layout

for video, poster, note in SEGMENTS:
    vpath = os.path.join(OUT, video)
    ppath = os.path.join(OUT, poster)
    if not os.path.exists(vpath):
        raise SystemExit(f"missing video: {vpath}")
    slide = prs.slides.add_slide(blank)
    # Full-bleed video with the rendered frame as its poster image.
    slide.shapes.add_movie(
        vpath,
        left=0, top=0,
        width=prs.slide_width, height=prs.slide_height,
        poster_frame_image=ppath if os.path.exists(ppath) else None,
        mime_type="video/mp4",
    )
    if AUTOPLAY:
        # python-pptx sets the media trigger to delay="indefinite" (wait for a
        # click). Flipping every timing condition to delay="0" makes the video
        # start automatically when the slide appears.
        timing = slide._element.find(qn("p:timing"))
        if timing is not None:
            for cond in timing.iter(qn("p:cond")):
                cond.set("delay", "0")
    # Speaker note = your voice-over cue for that beat.
    slide.notes_slide.notes_text_frame.text = note

suffix = "-autoplay" if AUTOPLAY else ""
dest = os.path.join(OUT, f"Fable5-presentation-dracula{suffix}.pptx")
prs.save(dest)
print("saved:", dest)
print("slides:", len(prs.slides._sldIdLst))
