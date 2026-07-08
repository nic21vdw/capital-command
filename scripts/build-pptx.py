"""Assemble the rendered Remotion segments into a PowerPoint deck where each
slide's video AUTO-PLAYS when the slide is shown in slideshow mode."""
import copy
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

OUT = "out"
BG = RGBColor(0x28, 0x2A, 0x36)  # Dracula background

# (video, poster still, slide title)
SLIDES = [
    ("titlecomp.mp4",        "stills/title.png",         "Manhattan Column Buckling, Explained"),
    ("headline.mp4",         "stills/headline_end.png",  "The headlines called it a beam"),
    ("columnvsbeamcomp.mp4", "stills/colvbeam.png",      "Column vs Beam"),
    ("buckling.mp4",         "stills/buckle_scurve.png", "Euler Buckling"),
    ("causescomp.mp4",       "stills/cause1.png",         "Four Likely Causes"),
]

prs = Presentation()
# 16:9 widescreen, matching the 1920x1080 video aspect exactly (full bleed).
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def make_autoplay(slide, movie):
    """Rewrite the slide timing so the media starts automatically (with the
    slide) instead of on click."""
    tree = slide._element
    timing = tree.find(qn("p:timing"))
    if timing is None:
        return
    # click-to-play -> start automatically
    for cond in timing.iter(qn("p:cond")):
        if cond.get("delay") == "indefinite":
            cond.set("delay", "0")
    for ctn in timing.iter(qn("p:cTn")):
        if ctn.get("nodeType") == "clickEffect":
            ctn.set("nodeType", "withEffect")


for video, poster, title in SLIDES:
    slide = prs.slides.add_slide(blank)

    # Dracula background fill.
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    # Full-bleed video with a poster frame so the static slide looks right.
    movie = slide.shapes.add_movie(
        f"{OUT}/{video}",
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
        poster_frame_image=f"{OUT}/{poster}",
        mime_type="video/mp4",
    )
    make_autoplay(slide, movie)

    # Speaker note naming the segment.
    slide.notes_slide.notes_text_frame.text = title

prs.save(f"{OUT}/manhattan-column-buckling.pptx")
print("wrote", f"{OUT}/manhattan-column-buckling.pptx")
